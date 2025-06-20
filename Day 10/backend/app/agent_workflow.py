import os
import json
import time
import uuid
from datetime import datetime
from dotenv import load_dotenv
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract

from langchain.chains import LLMChain
from langchain_core.prompts import PromptTemplate
from langchain.agents import Tool, initialize_agent
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_community.document_loaders import PyPDFLoader, PyPDFDirectoryLoader

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA
from langchain_google_genai import ChatGoogleGenerativeAI
from tavily import TavilyClient
import google.generativeai as genai

# ─────────────────────────────────────────────────────────────
# 2. API KEY SETUP
# ─────────────────────────────────────────────────────────────
load_dotenv()

os.environ["GOOGLE_API_KEY"] = os.getenv('GOOGLE_API_KEY')
os.environ["HUGGINGFACE_API_KEY"] = os.getenv('HUGGINGFACE_API_KEY')
os.environ["TAVILY_API_KEY"] = os.getenv('TAVILY_API_KEY')

genai.configure(api_key=os.environ["GOOGLE_API_KEY"])
llm = ChatGoogleGenerativeAI(
    model="gemini-1.5-flash", 
    temperature=0.2,
    convert_system_message_to_human=True)
tavily = TavilyClient(api_key=os.environ["TAVILY_API_KEY"])

# ─────────────────────────────────────────────────────────────
# 3. RAG PIPELINE SETUP
# ─────────────────────────────────────────────────────────────
os.makedirs("/content/uploads", exist_ok=True)

# pdf_files = {
#     "sample_pitch.pdf": "1P8RuIrT7aEjaeZ9tHcdYOQCGBWlvHhpf",
#     "uber_sample.pdf": "1QuV48dpItRKFgfOoGnrzPIBR34L71XrK"
# }

# for filename, file_id in pdf_files.items():
#     url = f"https://drive.google.com/uc?export=download&id={file_id}"
#     output_path = f"/content/papers/{filename}"
#     !wget -q --show-progress "{url}" -O "{output_path}"

loader = PyPDFDirectoryLoader("./app/uploads")
pages = loader.load()

splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
chunks = splitter.split_documents(pages)

embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
vectorstore = Chroma.from_documents(chunks, embedding_model, persist_directory="rag_db")

qa_chain = RetrievalQA.from_chain_type(
    llm=llm, 
    retriever=vectorstore.as_retriever(), 
    chain_type="stuff")

# ─────────────────────────────────────────────────────────────
# 4. TOOL DEFINITIONS
# ─────────────────────────────────────────────────────────────
VERSION_HISTORY = {}

tools = [
    Tool(
        name="SlideParserAgent",
        func=lambda file_text: LLMChain(
            llm=llm,
            prompt=PromptTemplate(
                input_variables=["file_text"],
                template=(
                    "You are a slide parser agent. Given pitch deck text:\n{file_text}\n"
                    "Segment it into structured JSON with fields like 'problem_slide', 'solution_slide', etc., "
                    "extract slide titles, main points, layout hints, and text content. Output JSON."
                    "also mention the slide count or slide number in the json data"
                )
            )
        ).run(file_text=file_text),
        description="Parse pitch deck text into structured JSON format."
    ),
    Tool(
        name="ContentEvaluatorAgent",
        func=lambda parsed_json: LLMChain(
            llm=llm,
            prompt=PromptTemplate(
                input_variables=["parsed_json"],
                template=(
                    "Evaluate each slide in the following parsed pitch deck JSON:\n{parsed_json}\n"
                    "Score each slide on Clarity, Narrative Coherence, and Investor Appeal (0-100), "
                    "Identify any weak areas and provide qualitative feedback."
                )
            )
        ).run(parsed_json=parsed_json),
        description="Evaluate slide content, flow, and investor appeal."
    ),
    Tool(
    name="ImprovementSuggesterAgent",
    func=lambda inputs: LLMChain(
        llm=llm,
        prompt=PromptTemplate(
            input_variables=["evaluation_json", "parsed_json"],
            template=(
                "You are an improvement suggester agent.\n\n"
                "Given the following pitch deck evaluation:\n{evaluation_json}\n\n"
                "And the parsed slide content:\n{parsed_json}\n\n"
                "For each slide, do the following:\n"
                "1. Suggest **two actionable improvements** (e.g., rewriting content, enhancing visual layout, or improving storytelling).\n"
                "2. Provide a **third point** explaining how the above two improvements will help align the slide better with **typical investor expectations**.\n\n"
                "Format your response slide-by-slide in a structured and clear format."
            )
        )
    ).run(evaluation_json=inputs["evaluation_json"], parsed_json=inputs["parsed_json"]),
    description="Generate 2 improvements and 1 investor-alignment rationale per slide."
),
    Tool(
        name="BenchmarkRetrieverAgent",
        func=lambda context_json: qa_chain.run(
            f"Show top pitch deck examples and best practices similar to: {context_json}"
        ),
        description="Retrieve benchmarked slide examples using RAG."
    ),
    Tool(
        name="VersionControlAgent",
        func=lambda deck_json: (
            lambda vid: (
                VERSION_HISTORY.setdefault(vid, []).append({
                    "id": str(uuid.uuid4()),
                    "timestamp": datetime.utcnow().isoformat(),
                    "deck": deck_json
                }),
                f"Saved version under ID: {vid}"
            )[1]
        )(str(hash(deck_json))),
        description="Store pitch deck version with timestamp for tracking evolution."
    )
]

# ─────────────────────────────────────────────────────────────
# 5. TEXT EXTRACTION FROM PDF / PPTX
# ─────────────────────────────────────────────────────────────
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        text.append("\n".join(slide_text))
    return "\n\n".join(text)

def extract_text_with_ocr(file_path):
    images = convert_from_path(file_path)
    return "\n\n".join([pytesseract.image_to_string(img) for img in images])

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        try:
            loader = PyPDFLoader(file_path)
            pages = loader.load()
            content = "\n".join([p.page_content for p in pages])
            return content if len(content.strip()) > 100 else extract_text_with_ocr(file_path)
        except:
            return extract_text_with_ocr(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format. Use PDF or PPTX.")

# ─────────────────────────────────────────────────────────────
# 6. RUN AGENTIC WORKFLOW
# ─────────────────────────────────────────────────────────────
agent = initialize_agent(tools, llm, agent="zero-shot-react-description", verbose=True)

def run_pitch_deck_workflow(file_path):
    deck_text = extract_text(file_path)
    parsed = agent.run(f"Parse this deck:\n{deck_text}")
    evaluated = agent.run(f"Evaluate this parsed deck:\n{parsed}")
    improved = agent.run(f"Suggest improvements:\n{evaluated}")
    benchmarked = agent.run(f"Benchmark similar decks:\n{parsed}")
    version = agent.run(f"Save version of this deck:\n{parsed}")

    return {
        "parsed": json.loads(parsed),
        "evaluated": evaluated,
        "improved": improved,
        "benchmarked": benchmarked,
        "version_log": VERSION_HISTORY
    }

# ____________________________________________________
#            Test the Workflow
# ____________________________________________________

file_path = "input_decks/aurasense_ai.pptx"
deck_text = extract_text(file_path)

# print("\n====================== Pitch Deck Content ======================")
# print(deck_text[:1500])


# print("\n====================== SlideParserAgent ======================")
parsed_output = tools[0].func(deck_text)
# print(parsed_output)

# print("\n====================== ContentEvaluatorAgent ======================")
evaluated_output = tools[1].func(parsed_output)
# print(evaluated_output)

# print("\n====================== ImprovementSuggesterAgent ======================")
improvement_output = tools[2].func({
    "evaluation_json": evaluated_output,
    "parsed_json": parsed_output})
# print(improvement_output)

# print("\n====================== BenchmarkRetrieverAgent ======================")
# benchmark_output = tools[3].func({
#     "parsed_json": parsed_output})
# print(benchmark_output)

# Step 4: Test RAG Chain
# query = "What problem does Airbnb solve?"
# response = qa_chain.run(query)
# print("Query:", query)
# print("Answer:", response)

print("\n====================== VersionControlAgent ======================")
version_output = tools[4].func(parsed_output)
print(version_output)

# # Show version history
# print("\n====================== Version History ======================")
# import pprint
# pprint.pprint(VERSION_HISTORY)